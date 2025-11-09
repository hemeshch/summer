from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

# Add blue gradient background effect
title_bg = slide_1.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(7.5))
title_frame = title_bg.text_frame
title_frame.clear()

# Main title
title_box = slide_1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
title_frame = title_box.text_frame
title_frame.clear()
p = title_frame.add_paragraph()
p.text = "The Fascinating World of Fish"
p.font.size = Pt(48)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(0, 51, 102)

# Subtitle
subtitle_box = slide_1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.clear()
p = subtitle_frame.add_paragraph()
p.text = "Exploring Aquatic Diversity and Adaptations"
p.font.size = Pt(24)
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(0, 102, 153)

# Add decorative element
deco_box = slide_1.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(0.5))
deco_frame = deco_box.text_frame
deco_frame.clear()
p = deco_frame.add_paragraph()
p.text = "🐠 🐟 🐡 🦈 🐙"
p.font.size = Pt(28)
p.alignment = PP_ALIGN.CENTER

# Slide 2: Introduction - What Makes a Fish?
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
title = slide_2.shapes.title
title.text = "What Makes a Fish?"

content = slide_2.placeholders[1]
tf = content.text_frame
tf.clear()

points = [
    ("Aquatic Vertebrates", [
        "Live in water throughout their lives",
        "Have a backbone (vertebrate)"
    ]),
    ("Key Characteristics", [
        "Breathe through gills",
        "Have fins for movement",
        "Most are cold-blooded (ectothermic)",
        "Bodies covered in scales (most species)"
    ]),
    ("Incredible Diversity", [
        "Over 34,000 known species",
        "Found in nearly every aquatic habitat",
        "Range from tiny gobies (8mm) to whale sharks (18m)"
    ])
]

for main_point, sub_points in points:
    p = tf.add_paragraph()
    p.text = main_point
    p.font.bold = True
    p.font.size = Pt(20)
    p.level = 0
    
    for sub_point in sub_points:
        p = tf.add_paragraph()
        p.text = sub_point
        p.font.size = Pt(16)
        p.level = 1

# Slide 3: Classification of Fish
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
title.text = "Major Groups of Fish"

content = slide_3.placeholders[1]
tf = content.text_frame
tf.clear()

groups = [
    ("Jawless Fish (Agnatha)", [
        "Most primitive group",
        "Examples: Lampreys and hagfish",
        "No paired fins or jaws"
    ]),
    ("Cartilaginous Fish (Chondrichthyes)", [
        "Skeleton made of cartilage",
        "Examples: Sharks, rays, and skates",
        "About 1,000 species"
    ]),
    ("Bony Fish (Osteichthyes)", [
        "Skeleton made of bone",
        "Largest group - over 30,000 species",
        "Includes most familiar fish: salmon, tuna, goldfish"
    ])
]

for group_name, details in groups:
    p = tf.add_paragraph()
    p.text = group_name
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.level = 0
    
    for detail in details:
        p = tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(16)
        p.level = 1

# Slide 4: Anatomy and Adaptations
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
title.text = "Fish Anatomy & Adaptations"

content = slide_4.placeholders[1]
tf = content.text_frame
tf.clear()

anatomy = [
    ("Gills", "Extract oxygen from water - incredibly efficient design"),
    ("Swim Bladder", "Gas-filled organ for buoyancy control (bony fish)"),
    ("Lateral Line", "Sensory system detecting water movement and pressure"),
    ("Fins", "Different types for different functions:"),
    ("", "• Dorsal & anal fins: Stability"),
    ("", "• Pectoral & pelvic fins: Steering and braking"),
    ("", "• Caudal fin (tail): Main propulsion"),
    ("Scales", "Protection and hydrodynamics - like natural armor"),
    ("Specialized Features", "Electric organs, bioluminescence, camouflage")
]

for feature, description in anatomy:
    p = tf.add_paragraph()
    if feature:
        p.text = f"{feature}: {description}" if description and not description.startswith("•") else feature
        p.font.bold = feature != ""
    else:
        p.text = description
    p.font.size = Pt(16)
    if description.startswith("•"):
        p.level = 1

# Slide 5: Habitats
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Diverse Aquatic Habitats"

content = slide_5.placeholders[1]
tf = content.text_frame
tf.clear()

habitats = [
    ("Freshwater Environments", [
        "Rivers and streams",
        "Lakes and ponds",
        "Wetlands and swamps",
        "Underground caves"
    ]),
    ("Marine Environments", [
        "Coral reefs - the rainforests of the sea",
        "Open ocean (pelagic zone)",
        "Deep sea - down to 11,000 meters",
        "Coastal waters and estuaries"
    ]),
    ("Extreme Environments", [
        "Arctic waters (antifreeze proteins)",
        "Hot springs (up to 44°C)",
        "Oxygen-poor waters",
        "Temporary pools (lungfish)"
    ])
]

for habitat_type, locations in habitats:
    p = tf.add_paragraph()
    p.text = habitat_type
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.level = 0
    
    for location in locations:
        p = tf.add_paragraph()
        p.text = location
        p.font.size = Pt(16)
        p.level = 1

# Slide 6: Fascinating Species
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
title.text = "Remarkable Fish Species"

content = slide_6.placeholders[1]
tf = content.text_frame
tf.clear()

species = [
    ("🦈 Great White Shark", "Apex predator with 300 teeth, can detect blood from miles away"),
    ("🐠 Clownfish", "Immune to sea anemone stings, all born male, can change to female"),
    ("🐡 Pufferfish", "Can inflate to 3x normal size, contains deadly tetrodotoxin"),
    ("⚡ Electric Eel", "Can generate 600 volts - enough to stun a horse"),
    ("🌟 Anglerfish", "Uses bioluminescent lure, extreme sexual dimorphism"),
    ("🐟 Coelacanth", "Living fossil - thought extinct for 66 million years"),
    ("🎨 Mandarin Fish", "No scales, covered in toxic mucus, spectacular colors"),
    ("🏔️ Antarctic Icefish", "No red blood cells - clear blood with antifreeze")
]

for name, fact in species:
    p = tf.add_paragraph()
    p.text = f"{name}"
    p.font.bold = True
    p.font.size = Pt(18)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = fact
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1

# Slide 7: Fish Behavior
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
title.text = "Fascinating Fish Behaviors"

content = slide_7.placeholders[1]
tf = content.text_frame
tf.clear()

behaviors = [
    ("Schooling", "Thousands move as one - protection and efficiency"),
    ("Migration", "Salmon travel thousands of miles to spawn"),
    ("Communication", "Through colors, sounds, electrical signals, and chemicals"),
    ("Parental Care", "Seahorses: males carry eggs; Cichlids: mouth brooders"),
    ("Hunting Strategies", "Archerfish shoot water jets; Anglerfish use lures"),
    ("Symbiosis", "Cleaner fish provide spa services to larger fish"),
    ("Tool Use", "Some wrasses use rocks to crack open shells"),
    ("Sleep Patterns", "Parrotfish create mucus sleeping bags")
]

for behavior, description in behaviors:
    p = tf.add_paragraph()
    p.text = f"• {behavior}: {description}"
    p.font.size = Pt(16)
    if ":" in p.text:
        p.font.color.rgb = RGBColor(0, 51, 102)

# Slide 8: Ecological Importance
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_8.shapes.title
title.text = "Why Fish Matter to Our Planet"

content = slide_8.placeholders[1]
tf = content.text_frame
tf.clear()

importance = [
    ("Food Web Foundation", [
        "Key link in aquatic food chains",
        "Feed billions of people worldwide",
        "Support entire ecosystems"
    ]),
    ("Environmental Services", [
        "Nutrient cycling in water bodies",
        "Coral reef maintenance (herbivorous fish)",
        "Seed dispersal in flooded forests",
        "Bioturbation - mixing sediments"
    ]),
    ("Economic Value", [
        "Commercial fishing: $150+ billion annually",
        "Recreational fishing and tourism",
        "Aquarium trade and ornamental fish"
    ]),
    ("Scientific & Medical Benefits", [
        "Model organisms for research",
        "Source of omega-3 fatty acids",
        "Potential medicines from fish compounds"
    ])
]

for category, points in importance:
    p = tf.add_paragraph()
    p.text = category
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 102, 51)
    p.level = 0
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.level = 1

# Slide 9: Conservation Challenges
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_9.shapes.title
title.text = "Conservation Challenges & Solutions"

content = slide_9.placeholders[1]
tf = content.text_frame
tf.clear()

conservation = [
    ("Major Threats", [
        "Overfishing and bycatch",
        "Habitat destruction (coral reefs, mangroves)",
        "Pollution (plastics, chemicals, noise)",
        "Climate change and ocean acidification",
        "Invasive species"
    ]),
    ("Conservation Efforts", [
        "Marine Protected Areas (MPAs)",
        "Sustainable fishing practices",
        "Habitat restoration projects",
        "Captive breeding programs",
        "International agreements and regulations"
    ]),
    ("How You Can Help", [
        "Choose sustainable seafood",
        "Reduce plastic use",
        "Support conservation organizations",
        "Learn and share knowledge",
        "Participate in citizen science"
    ])
]

for section, items in conservation:
    p = tf.add_paragraph()
    p.text = section
    p.font.bold = True
    p.font.size = Pt(20)
    if "Threats" in section:
        p.font.color.rgb = RGBColor(153, 0, 0)
    elif "Efforts" in section:
        p.font.color.rgb = RGBColor(0, 102, 51)
    else:
        p.font.color.rgb = RGBColor(0, 51, 102)
    p.level = 0
    
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 1

# Slide 10: Amazing Fish Facts
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_10.shapes.title
title.text = "Amazing Fish Facts"

content = slide_10.placeholders[1]
tf = content.text_frame
tf.clear()

facts = [
    "🎯 Fish have been on Earth for over 500 million years",
    "👁️ Four-eyed fish (Anableps) can see above and below water simultaneously",
    "💡 Over 1,500 fish species can produce light (bioluminescence)",
    "🏃 Sailfish can swim up to 68 mph - fastest fish in the ocean",
    "🎨 Fish can see ultraviolet and polarized light invisible to humans",
    "💤 Some sharks must swim constantly to breathe",
    "🧠 Fish can recognize human faces and have demonstrated tool use",
    "🌡️ Polar fish have antifreeze proteins in their blood",
    "📏 The smallest fish (Paedocypris) is only 7.9mm long",
    "⏰ Koi fish can live over 200 years",
    "🗣️ Fish communicate through grunts, croaks, and drumming sounds"
]

for fact in facts:
    p = tf.add_paragraph()
    p.text = fact
    p.font.size = Pt(16)
    p.level = 0

# Slide 11: Conclusion
slide_11 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

# Title
title_box = slide_11.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.clear()
p = title_frame.add_paragraph()
p.text = "The Underwater World Awaits"
p.font.size = Pt(40)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(0, 51, 102)

# Main message
msg_box = slide_11.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(3))
msg_frame = msg_box.text_frame
msg_frame.clear()

messages = [
    "Fish represent one of nature's greatest success stories",
    "",
    "From the smallest goby to the mighty whale shark, each species plays a vital role in our planet's aquatic ecosystems",
    "",
    "By understanding and protecting fish, we preserve the health of our oceans, rivers, and lakes for future generations"
]

for msg in messages:
    p = msg_frame.add_paragraph()
    p.text = msg
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    if msg and "success" in msg:
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 51)
    elif msg and "vital role" in msg:
        p.font.color.rgb = RGBColor(0, 51, 102)
    elif msg and "future" in msg:
        p.font.italic = True
        p.font.color.rgb = RGBColor(0, 102, 153)

# Emoji decoration
deco_box = slide_11.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(0.8))
deco_frame = deco_box.text_frame
deco_frame.clear()
p = deco_frame.add_paragraph()
p.text = "🌊 🐠 🐟 🐡 🦈 🐙 🌊"
p.font.size = Pt(32)
p.alignment = PP_ALIGN.CENTER

# Thank you
thanks_box = slide_11.shapes.add_textbox(Inches(3), Inches(6.5), Inches(4), Inches(0.5))
thanks_frame = thanks_box.text_frame
thanks_frame.clear()
p = thanks_frame.add_paragraph()
p.text = "Thank You!"
p.font.size = Pt(24)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(0, 51, 102)

# Save presentation
prs.save('/mnt/conversation_data/fish_presentation.pptx')
print("Presentation created successfully!")
print("Total slides: 11")
print("Topics covered: Introduction, Classification, Anatomy, Habitats, Species, Behaviors, Ecology, Conservation, Facts")
