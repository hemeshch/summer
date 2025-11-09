from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Set slide size to widescreen (16:9)
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add a slide with Title and Content layout
slide_layout = prs.slide_layouts[5]  # Blank layout for more control
slide = prs.slides.add_slide(slide_layout)

# Add title
title_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
)
title_frame = title_box.text_frame
title_frame.clear()
p = title_frame.add_paragraph()
p.text = "Fish: Life Under Water"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 84, 164)  # Ocean blue color
p.alignment = PP_ALIGN.CENTER

# Add main content box with key facts
content_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(1.3), Inches(4.5), Inches(3.8)
)
content_frame = content_box.text_frame
content_frame.clear()

# Key Facts section
facts_title = content_frame.add_paragraph()
facts_title.text = "Key Facts About Fish"
facts_title.font.size = Pt(20)
facts_title.font.bold = True
facts_title.font.color.rgb = RGBColor(0, 0, 0)
facts_title.space_after = Pt(6)

# Add facts as bullet points
facts = [
    "Over 34,000 species worldwide",
    "Breathe through gills",
    "Cold-blooded vertebrates",
    "Live in freshwater and saltwater",
    "First appeared 530 million years ago"
]

for fact in facts:
    p = content_frame.add_paragraph()
    p.text = "• " + fact
    p.font.size = Pt(14)
    p.level = 0
    p.space_after = Pt(8)

# Add interesting features box
features_box = slide.shapes.add_textbox(
    Inches(5.2), Inches(1.3), Inches(4.3), Inches(3.8)
)
features_frame = features_box.text_frame
features_frame.clear()

# Amazing Features section
features_title = features_frame.add_paragraph()
features_title.text = "Amazing Features"
features_title.font.size = Pt(20)
features_title.font.bold = True
features_title.font.color.rgb = RGBColor(0, 0, 0)
features_title.space_after = Pt(6)

# Add features
features = [
    "Lateral line system detects vibrations",
    "Some fish can change colors",
    "Sharks have cartilage instead of bones",
    "Many species migrate thousands of miles",
    "Some deep-sea fish produce light"
]

for feature in features:
    p = features_frame.add_paragraph()
    p.text = "• " + feature
    p.font.size = Pt(14)
    p.level = 0
    p.space_after = Pt(8)

# Add a footer with a fun fact
footer_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(4.8), Inches(9), Inches(0.5)
)
footer_frame = footer_box.text_frame
footer_frame.clear()
p = footer_frame.add_paragraph()
p.text = "🐠 Fun Fact: The largest fish is the whale shark, reaching up to 40 feet long! 🐠"
p.font.size = Pt(12)
p.font.italic = True
p.font.color.rgb = RGBColor(0, 84, 164)
p.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('/mnt/conversation_data/fish_presentation.pptx')
print("PowerPoint presentation created successfully!")