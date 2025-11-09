from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
TITLE_COLOR = RGBColor(44, 62, 80)  # Dark blue-grey
ACCENT_COLOR = RGBColor(52, 152, 219)  # Bright blue
TEXT_COLOR = RGBColor(52, 73, 94)  # Dark grey
LIGHT_BG = RGBColor(236, 240, 241)  # Light grey

def add_title_slide():
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Stella Technical Stack"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Architecture, Capabilities & Implementation"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Created by Jess Solutions Incorporated"
    footer_frame.paragraphs[0].font.size = Pt(14)
    footer_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_overview_slide():
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Overview"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Stella is a productivity-focused AI assistant built on advanced language model architecture"
    tf.paragraphs[0].font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Based on Claude architecture by Anthropic"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Specialized for computer control and automation"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Operates in sandboxed Linux container environment"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Integrates with host machine for specific tasks"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Knowledge cutoff: January 2025"
    p.font.size = Pt(16)
    p.level = 0

def add_core_architecture_slide():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Core Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Foundation Model"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Claude-based large language model"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Advanced natural language understanding"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nExecution Environment"
    p.font.size = Pt(20)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Docker containerized Linux (Ubuntu 24.04)"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Isolated workspace at /mnt directory"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Persistent container per conversation"
    p.font.size = Pt(16)
    p.level = 1

def add_tool_categories_slide():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Tool Categories"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Container Tools"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• ZSH command execution"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• File management & editing"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nHost Machine Tools"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Calendar integration (Apple Calendar)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Email sending (Apple Mail)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Codebase analysis (Claude Code)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nWeb & Information Tools"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Web search and content fetching"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Image viewing and analysis"
    p.font.size = Pt(14)
    p.level = 1

def add_container_stack_slide():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Container Technology Stack"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Operating System"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Ubuntu 24.04 LTS with ZSH shell"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nProgramming Languages"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Python 3.12 (NumPy, Pandas, Matplotlib, python-pptx)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Node.js 18.x (pptxgenjs, markitdown)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nDocument Processing"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Pandoc for format conversion"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• LibreOffice for document manipulation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• ImageMagick for image processing"
    p.font.size = Pt(14)
    p.level = 1

def add_capabilities_slide():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Key Capabilities"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Development & Coding"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Write, execute, and debug code"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Analyze existing codebases"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nDocument Creation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Generate presentations (PowerPoint)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Create reports and documentation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nProductivity Management"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Calendar event retrieval"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Email composition and sending"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nInformation Processing"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Web search and content analysis"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Image viewing and interpretation"
    p.font.size = Pt(14)
    p.level = 1

def add_security_slide():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Security & Isolation"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Container Isolation"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Docker containerization for process isolation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Separate workspace per conversation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• No direct host filesystem access"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nControlled Host Access"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Limited to specific approved tools"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Read-only calendar access"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Explicit user authorization for emails"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nData Protection"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Copyright compliance in web searches"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• User fact storage with consent"
    p.font.size = Pt(14)
    p.level = 1

def add_integration_slide():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "System Integration"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "File Transfer"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• conversation_data/ folder for user downloads"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Direct file submission to users"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nTool Orchestration"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Function calling through XML protocol"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Parallel tool execution capability"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Error handling and recovery"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nMemory & Context"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• User fact remembering system"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Conversation context maintenance"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Citation and source tracking"
    p.font.size = Pt(14)
    p.level = 1

def add_use_cases_slide():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Primary Use Cases"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Software Development"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Code generation and debugging"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Project structure analysis"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nDocument Automation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Presentation creation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Report generation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nProductivity Enhancement"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Task automation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Email and calendar management"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nResearch & Analysis"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Web research and summarization"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Data processing and visualization"
    p.font.size = Pt(14)
    p.level = 1

def add_summary_slide():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Summary"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Stella combines:"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "\n• Advanced language understanding (Claude architecture)"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "\n• Secure containerized execution environment"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "\n• Rich tool ecosystem for productivity"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "\n• Seamless integration with host systems"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "\n• Focus on practical task automation"
    p.font.size = Pt(16)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "\n\nDesigned to enhance productivity through intelligent automation"
    p.font.size = Pt(18)
    p.font.italic = True
    p.level = 0

# Create all slides
add_title_slide()
add_overview_slide()
add_core_architecture_slide()
add_tool_categories_slide()
add_container_stack_slide()
add_capabilities_slide()
add_security_slide()
add_integration_slide()
add_use_cases_slide()
add_summary_slide()

# Save the presentation
prs.save('conversation_data/stella_technical_stack.pptx')
print("Presentation created successfully: stella_technical_stack.pptx")
