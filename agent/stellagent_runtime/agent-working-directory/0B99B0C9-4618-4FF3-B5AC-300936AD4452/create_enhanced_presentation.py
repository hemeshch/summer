from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR
import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
TITLE_COLOR = RGBColor(44, 62, 80)  # Dark blue-grey
ACCENT_COLOR = RGBColor(52, 152, 219)  # Bright blue
TEXT_COLOR = RGBColor(52, 73, 94)  # Dark grey
HIGHLIGHT_COLOR = RGBColor(46, 204, 113)  # Green for emphasis
CODE_BG = RGBColor(241, 243, 245)  # Light grey for code

def add_title_slide():
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Stella Technical Deep Dive"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Production Architecture & Implementation Details\nBased on Actual Codebase Analysis"
    for p in subtitle_frame.paragraphs:
        p.font.size = Pt(22)
        p.font.color.rgb = ACCENT_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    # Add footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Created by Jess Solutions Incorporated"
    footer_frame.paragraphs[0].font.size = Pt(14)
    footer_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_architecture_overview():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Multi-Platform Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Three-Tier Distributed System"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\nPython Client Layer"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• AsyncIO-based architecture with WebSocket client"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Container orchestration via Docker SDK"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nCloudflare Workers Server"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Edge computing with Durable Objects for state"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Real-time WebSocket broadcasting"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\niOS SwiftUI Application"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Native speech-to-text integration"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Persistent conversation storage"
    p.font.size = Pt(14)
    p.level = 1

def add_container_innovation():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Revolutionary Container Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Per-Conversation Isolation"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\nEach conversation gets dedicated Ubuntu 24.04 container:"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "• 4GB RAM limit, 4 CPU cores"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Persistent /mnt working directory"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Pre-installed development stack:"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  - Python 3.12 with 50+ data science packages"
    p.font.size = Pt(13)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "  - Node.js 18.x with document tools"
    p.font.size = Pt(13)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "  - PowerPoint creation (pptxgenjs, python-pptx)"
    p.font.size = Pt(13)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "  - Document processing (Pandoc, LibreOffice)"
    p.font.size = Pt(13)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "\nLifecycle Management:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Automatic cleanup after inactivity"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Container reset with file preservation option"
    p.font.size = Pt(14)
    p.level = 1

def add_claude_integration():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Sophisticated Claude API Integration"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Production-Ready Implementation"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\nInfinite Retry Logic"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Exponential backoff (1s → 2s → 4s → 8s)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Automatic model fallback:"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "  claude-opus-4 → claude-3-haiku on failure"
    p.font.size = Pt(13)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "\nReal-Time Streaming"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• SSE (Server-Sent Events) for responses"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Early tool call detection and broadcasting"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Multi-tool execution in single response"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nAdvanced Features"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• System prompt injection with user memory"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Token usage tracking and optimization"
    p.font.size = Pt(14)
    p.level = 1

def add_tool_system():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Plugin-Based Tool System"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "BaseToolSetProvider Architecture"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\nCore Tool Categories:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Container Tools: ZSH, file editing, status management"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Host Tools: Calendar, email, codebase analysis"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Web Tools: Search, fetch, content analysis"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Memory Tools: User fact storage"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nState Management:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Global state + per-conversation state"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Atomic updates with validation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Pickle serialization for persistence"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nExtensibility:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Abstract base class for new providers"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• JSON schema validation"
    p.font.size = Pt(14)
    p.level = 1

def add_realtime_features():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Real-Time Communication System"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "WebSocket Broadcasting Architecture"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\nLive Tool Execution Visibility:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Tool calls broadcast as they start"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Human-readable tool names"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Status updates (running, completed, failed)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nCloudflare Durable Objects:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Persistent WebSocket connections"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Conversation state management"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Multi-client synchronization"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nConnection Resilience:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Automatic reconnection with backoff"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Message queue during disconnection"
    p.font.size = Pt(14)
    p.level = 1

def add_memory_system():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Cross-Conversation AI Memory"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Persistent User Context System"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\nImplementation:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• File-based storage with atomic writes"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Facts injected into system prompts"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Per-user isolation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nMemory Types:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• User preferences and settings"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Project context and paths"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Personal information (with consent)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nIntegration:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Automatic loading on conversation start"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Real-time updates during conversation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Explicit user control (remember/forget)"
    p.font.size = Pt(14)
    p.level = 1

def add_security_features():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Security & Error Handling"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Production-Grade Security"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\nContainer Isolation:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Resource limits (RAM, CPU, disk)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Network isolation per container"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Path traversal protection"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nError Recovery:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Automatic model switching on failure"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Container health monitoring"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Graceful degradation strategies"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nData Protection:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Conversation isolation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Secure file handling"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• API key management"
    p.font.size = Pt(14)
    p.level = 1

def add_performance_optimizations():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Performance Optimizations"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Production Performance Features"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\nCaching Strategy:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Multi-level state caching"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Container reuse for same conversation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Tool result caching"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nAsynchronous Processing:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• AsyncIO for concurrent operations"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Parallel tool execution"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Non-blocking WebSocket I/O"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nResource Management:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Automatic container cleanup"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Memory-efficient streaming"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Connection pooling"
    p.font.size = Pt(14)
    p.level = 1

def add_technical_stack_detail():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Complete Technical Stack"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Backend (Python 3.8+)"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    p = tf.add_paragraph()
    p.text = "• AsyncIO, aiohttp, websockets"
    p.font.size = Pt(13)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Docker SDK, container management"
    p.font.size = Pt(13)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 100+ packages for data science & docs"
    p.font.size = Pt(13)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nServer (Cloudflare Workers)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    p = tf.add_paragraph()
    p.text = "• Edge computing runtime"
    p.font.size = Pt(13)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Durable Objects for state"
    p.font.size = Pt(13)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• WebSocket API gateway"
    p.font.size = Pt(13)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nMobile (iOS SwiftUI)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    p = tf.add_paragraph()
    p.text = "• Native UI components"
    p.font.size = Pt(13)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Speech recognition"
    p.font.size = Pt(13)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Core Data persistence"
    p.font.size = Pt(13)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nIntegrations"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    p = tf.add_paragraph()
    p.text = "• Anthropic Claude API"
    p.font.size = Pt(13)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Apple Shortcuts (Calendar/Mail)"
    p.font.size = Pt(13)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Web APIs for search & fetch"
    p.font.size = Pt(13)
    p.level = 1

def add_unique_features():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Unique Innovation Highlights"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "What Makes Stella Special"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\n✨ Container-Per-Conversation Model"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Industry-first isolated development environment"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n✨ Real-Time Tool Visibility"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Users see exactly what's happening as it happens"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n✨ Infinite Retry with Fallback"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Never fails - automatically switches models"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n✨ Cross-Platform Memory"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Remembers user context across all devices"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n✨ Production-Ready Architecture"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Enterprise-grade error handling & security"
    p.font.size = Pt(14)
    p.level = 1

def add_business_value():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Business Value Proposition"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Enterprise Productivity Solution"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\nKey Differentiators:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Full development environment per conversation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Real-time transparency of AI actions"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Never-fail architecture with automatic recovery"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Cross-platform synchronization"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nUse Case Coverage:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Software development & debugging"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Document & presentation creation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Data analysis & visualization"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Email & calendar automation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Research & content generation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nTarget Market:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Enterprise teams needing AI automation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Professional developers & data scientists"
    p.font.size = Pt(14)
    p.level = 1

def add_future_roadmap():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Technical Roadmap"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Future Enhancements"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\nShort Term (Q1 2025):"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Additional tool integrations"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Enhanced memory capabilities"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Performance optimizations"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nMedium Term (Q2-Q3 2025):"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Multi-agent collaboration"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Advanced workflow automation"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Enterprise security features"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nLong Term (2026+):"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Custom model fine-tuning"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Autonomous task execution"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Industry-specific solutions"
    p.font.size = Pt(14)
    p.level = 1

def add_summary():
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Stella: Production-Ready AI Assistant"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = HIGHLIGHT_COLOR
    
    p = tf.add_paragraph()
    p.text = "\n✅ Multi-platform architecture (Python + Cloudflare + iOS)"
    p.font.size = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "\n✅ Revolutionary container-per-conversation model"
    p.font.size = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "\n✅ Sophisticated Claude integration with infinite retry"
    p.font.size = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "\n✅ Real-time WebSocket communication"
    p.font.size = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "\n✅ Cross-conversation memory system"
    p.font.size = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "\n✅ Production-grade security & error handling"
    p.font.size = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "\n✅ Extensible plugin architecture"
    p.font.size = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "\n\n🚀 Ready for enterprise deployment"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.italic = True

# Create all slides
add_title_slide()
add_architecture_overview()
add_container_innovation()
add_claude_integration()
add_tool_system()
add_realtime_features()
add_memory_system()
add_security_features()
add_performance_optimizations()
add_technical_stack_detail()
add_unique_features()
add_business_value()
add_future_roadmap()
add_summary()

# Save the presentation
prs.save('conversation_data/stella_technical_deep_dive.pptx')
print("Enhanced presentation created successfully!")
print("Total slides: 14")
print("File saved as: stella_technical_deep_dive.pptx")
