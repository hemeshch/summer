from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

def create_stella_architecture_presentation():
    # Create presentation
    prs = Presentation()
    
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Stella Technical Architecture"
    subtitle.text = "AI Assistant System Design & Implementation\nBuilt on Claude Architecture by Anthropic"
    
    # Slide 2: Overview
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "System Overview"
    content.text = ("• AI assistant focused on productivity and task automation\n"
                   "• Built on Claude architecture (Anthropic foundation)\n"
                   "• Developed by Jess Solutions Incorporated\n"
                   "• Operates through sandboxed container environment\n"
                   "• Integrates with host system for specific operations\n"
                   "• Natural language processing with tool-use capabilities")
    
    # Slide 3: Core Architecture Components
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Core Architecture Components"
    content.text = ("1. Language Model Core\n"
                   "   • Claude-based large language model\n"
                   "   • Advanced natural language understanding\n"
                   "   • Context-aware response generation\n\n"
                   "2. Tool Integration Layer\n"
                   "   • Function calling interface\n"
                   "   • Parameter validation and execution\n\n"
                   "3. Execution Environment\n"
                   "   • Docker containerization\n"
                   "   • Isolated Linux workspace")
    
    # Slide 4: Container Environment
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Sandboxed Container Environment"
    content.text = ("Technical Specifications:\n"
                   "• Ubuntu 24.04 LTS base image\n"
                   "• ZSH shell for command execution\n"
                   "• Python 3.12 with data science libraries\n"
                   "• Node.js 18.x for JavaScript operations\n"
                   "• Persistent /mnt directory across session\n\n"
                   "Pre-installed Tools:\n"
                   "• Document processing: pandoc, LibreOffice\n"
                   "• PowerPoint creation: python-pptx, pptxgenjs\n"
                   "• Image manipulation: ImageMagick\n"
                   "• Development tools: git, compilers")
    
    # Slide 5: Tool Capabilities
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Integrated Tool Capabilities"
    content.text = ("File Management:\n"
                   "• View, create, edit text files\n"
                   "• Navigate directory structures\n"
                   "• Image viewing and processing\n\n"
                   "System Operations:\n"
                   "• Execute shell commands (ZSH)\n"
                   "• Container status monitoring\n"
                   "• Environment reset with file preservation\n\n"
                   "External Integration:\n"
                   "• Web search and content fetching\n"
                   "• Email sending (Apple Mail)\n"
                   "• Calendar access (Apple Calendar)\n"
                   "• Codebase analysis (Claude Code)")
    
    # Slide 6: Security Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Security & Isolation"
    content.text = ("Container Isolation:\n"
                   "• Sandboxed Docker environment\n"
                   "• No direct host filesystem access\n"
                   "• Resource limits and constraints\n"
                   "• Separate workspace per conversation\n\n"
                   "Controlled Host Access:\n"
                   "• Specific tools for host operations\n"
                   "• Explicit permission boundaries\n"
                   "• Read-only codebase analysis\n\n"
                   "Data Protection:\n"
                   "• Session-based isolation\n"
                   "• No cross-conversation data access")
    
    # Slide 7: Processing Pipeline
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Request Processing Pipeline"
    content.text = ("1. Input Processing\n"
                   "   • Natural language parsing\n"
                   "   • Intent recognition\n"
                   "   • Context integration\n\n"
                   "2. Planning & Reasoning\n"
                   "   • Task decomposition\n"
                   "   • Tool selection logic\n"
                   "   • Parameter extraction\n\n"
                   "3. Execution\n"
                   "   • Function invocation\n"
                   "   • Result processing\n"
                   "   • Error handling\n\n"
                   "4. Response Generation\n"
                   "   • Result synthesis\n"
                   "   • Natural language formatting")
    
    # Slide 8: Memory & Context
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Memory & Context Management"
    content.text = ("Conversation Context:\n"
                   "• Full conversation history retention\n"
                   "• Tool execution results tracking\n"
                   "• Error and retry management\n\n"
                   "User Memory System:\n"
                   "• Persistent fact storage\n"
                   "• Cross-conversation memory\n"
                   "• User preference tracking\n\n"
                   "Knowledge Base:\n"
                   "• Training data up to January 2025\n"
                   "• Real-time web search integration\n"
                   "• Document analysis capabilities")
    
    # Slide 9: Technical Stack
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Stack Details"
    content.text = ("Programming Languages:\n"
                   "• Python 3.12 - Primary scripting\n"
                   "• JavaScript/Node.js - Web tools\n"
                   "• Shell scripting - System operations\n\n"
                   "Key Libraries:\n"
                   "• python-pptx - Presentation creation\n"
                   "• pandas - Data manipulation\n"
                   "• Pillow - Image processing\n"
                   "• markitdown - Document conversion\n\n"
                   "Infrastructure:\n"
                   "• Docker containerization\n"
                   "• Linux-based execution environment\n"
                   "• RESTful API interfaces")
    
    # Slide 10: Performance Characteristics
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Performance Characteristics"
    content.text = ("Response Generation:\n"
                   "• Context-aware processing\n"
                   "• Streaming response capability\n"
                   "• Multi-turn conversation support\n\n"
                   "Execution Limits:\n"
                   "• 300-second default command timeout\n"
                   "• 10-minute codebase analysis timeout\n"
                   "• File size handling up to 10MB\n\n"
                   "Optimization Features:\n"
                   "• Parallel tool execution\n"
                   "• Result caching within session\n"
                   "• Efficient memory management")
    
    # Slide 11: Use Cases
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Primary Use Cases"
    content.text = ("Document Creation:\n"
                   "• PowerPoint presentations\n"
                   "• Technical reports\n"
                   "• Data analysis documents\n\n"
                   "Development Support:\n"
                   "• Code writing and debugging\n"
                   "• Codebase analysis\n"
                   "• Script automation\n\n"
                   "Productivity Tasks:\n"
                   "• Email composition\n"
                   "• Calendar management\n"
                   "• File organization\n\n"
                   "Information Processing:\n"
                   "• Web research\n"
                   "• Data extraction and analysis")
    
    # Slide 12: Future Enhancements
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Architecture Advantages"
    content.text = ("Flexibility:\n"
                   "• Extensible tool framework\n"
                   "• Multiple language support\n"
                   "• Adaptable to various workflows\n\n"
                   "Reliability:\n"
                   "• Isolated execution environment\n"
                   "• Error recovery mechanisms\n"
                   "• Consistent performance\n\n"
                   "Scalability:\n"
                   "• Container-based architecture\n"
                   "• Modular component design\n"
                   "• Efficient resource utilization\n\n"
                   "User Experience:\n"
                   "• Natural language interface\n"
                   "• Context-aware responses\n"
                   "• Seamless tool integration")
    
    # Slide 13: Summary
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Summary"
    content.text = ("Stella represents a sophisticated AI assistant architecture that:\n\n"
                   "• Combines advanced language understanding with practical tool execution\n\n"
                   "• Maintains security through containerized isolation\n\n"
                   "• Provides comprehensive productivity capabilities\n\n"
                   "• Integrates seamlessly with user workflows\n\n"
                   "• Delivers consistent, reliable performance\n\n"
                   "Built on Claude's foundation with custom enhancements by Jess Solutions")
    
    # Save presentation
    prs.save('conversation_data/stella_architecture.pptx')
    print("Presentation created successfully!")

# Create the presentation
create_stella_architecture_presentation()
