from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

def add_rounded_rectangle(slide, left, top, width, height, color_rgb):
    """Add a rounded rectangle shape with custom color"""
    shape = slide.shapes.add_shape(
        1,  # Rounded rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color_rgb
    shape.line.fill.background()
    return shape

def style_title_slide(slide):
    """Apply custom styling to title slide"""
    # Add gradient background effect with shapes
    bg_shape = add_rounded_rectangle(
        slide, 
        Inches(0), Inches(0), 
        Inches(10), Inches(5.625),
        RGBColor(25, 31, 44)  # Dark blue background
    )
    bg_shape.shadow.visible = False
    
    # Move to back
    slide.shapes._spTree.remove(bg_shape._element)
    slide.shapes._spTree.insert(2, bg_shape._element)
    
    # Style title
    title = slide.shapes.title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Style subtitle
    subtitle = slide.placeholders[1]
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(147, 197, 253)  # Light blue
        paragraph.alignment = PP_ALIGN.CENTER

def style_content_slide(slide, accent_color=None):
    """Apply custom styling to content slides"""
    if accent_color is None:
        accent_color = RGBColor(59, 130, 246)  # Default blue
    
    # Add subtle background
    bg_shape = add_rounded_rectangle(
        slide,
        Inches(0), Inches(0),
        Inches(10), Inches(5.625),
        RGBColor(248, 250, 252)  # Very light gray
    )
    slide.shapes._spTree.remove(bg_shape._element)
    slide.shapes._spTree.insert(2, bg_shape._element)
    
    # Add accent bar at top
    accent_bar = add_rounded_rectangle(
        slide,
        Inches(0), Inches(0),
        Inches(10), Inches(0.08),
        accent_color
    )
    
    # Style title if it exists
    if slide.shapes.title:
        title = slide.shapes.title
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(25, 31, 44)
    
    # Style content if it exists
    if len(slide.placeholders) > 1:
        content = slide.placeholders[1]
        content.text_frame.margin_left = Inches(0.3)
        content.text_frame.margin_right = Inches(0.3)
        content.text_frame.margin_top = Inches(0.2)
        content.text_frame.margin_bottom = Inches(0.2)

def style_blank_slide(slide, accent_color=None):
    """Apply custom styling to blank slides"""
    if accent_color is None:
        accent_color = RGBColor(59, 130, 246)
    
    # Add subtle background
    bg_shape = add_rounded_rectangle(
        slide,
        Inches(0), Inches(0),
        Inches(10), Inches(5.625),
        RGBColor(248, 250, 252)
    )
    slide.shapes._spTree.remove(bg_shape._element)
    slide.shapes._spTree.insert(2, bg_shape._element)
    
    # Add accent bar at top
    accent_bar = add_rounded_rectangle(
        slide,
        Inches(0), Inches(0),
        Inches(10), Inches(0.08),
        accent_color
    )

def create_aesthetic_stella_presentation():
    # Create presentation
    prs = Presentation()
    
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Define color scheme
    primary_blue = RGBColor(59, 130, 246)
    secondary_blue = RGBColor(147, 197, 253)
    dark_blue = RGBColor(25, 31, 44)
    accent_purple = RGBColor(139, 92, 246)
    accent_teal = RGBColor(20, 184, 166)
    text_dark = RGBColor(31, 41, 55)
    text_light = RGBColor(107, 114, 128)
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    style_title_slide(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Stella\nTechnical Architecture"
    subtitle.text = "AI Assistant System Design & Implementation\nPowered by Claude • Developed by Jess Solutions"
    
    # Add decorative elements
    deco1 = add_rounded_rectangle(
        slide,
        Inches(0.5), Inches(3.5),
        Inches(0.3), Inches(0.3),
        secondary_blue
    )
    deco2 = add_rounded_rectangle(
        slide,
        Inches(9.2), Inches(1.5),
        Inches(0.3), Inches(0.3),
        accent_purple
    )
    
    # Slide 2: Overview with icon-like elements
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_content_slide(slide, primary_blue)
    
    title = slide.shapes.title
    title.text = "🚀 System Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    points = [
        ("🤖", "AI assistant focused on productivity and task automation"),
        ("🏗️", "Built on Claude architecture (Anthropic foundation)"),
        ("💼", "Developed by Jess Solutions Incorporated"),
        ("🔒", "Operates through sandboxed container environment"),
        ("🔗", "Integrates with host system for specific operations"),
        ("💬", "Natural language processing with tool-use capabilities")
    ]
    
    for icon, text in points:
        p = tf.add_paragraph()
        p.text = f"{icon}  {text}"
        p.font.size = Pt(16)
        p.font.color.rgb = text_dark
        p.space_after = Pt(12)
        p.level = 0
    
    # Slide 3: Core Architecture - Visual blocks
    slide_layout = prs.slide_layouts[5]  # Blank layout for custom design
    slide = prs.slides.add_slide(slide_layout)
    style_blank_slide(slide, accent_purple)
    
    # Add title manually
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "⚙️ Core Architecture Components"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = dark_blue
    
    # Create three component boxes
    components = [
        ("Language Model Core", "• Claude-based LLM\n• Natural language understanding\n• Context-aware responses", primary_blue),
        ("Tool Integration", "• Function calling interface\n• Parameter validation\n• Execution management", accent_purple),
        ("Execution Environment", "• Docker containerization\n• Isolated Linux workspace\n• Resource management", accent_teal)
    ]
    
    for i, (comp_title, comp_text, color) in enumerate(components):
        # Component box
        left = Inches(0.5 + i * 3.2)
        box = add_rounded_rectangle(
            slide,
            left, Inches(1.5),
            Inches(2.9), Inches(3.2),
            RGBColor(255, 255, 255)
        )
        box.shadow.visible = True
        box.shadow.blur_radius = Pt(10)
        box.shadow.distance = Pt(3)
        box.shadow.angle = 90
        
        # Component header
        header = add_rounded_rectangle(
            slide,
            left, Inches(1.5),
            Inches(2.9), Inches(0.6),
            color
        )
        
        # Add text
        text_box = slide.shapes.add_textbox(
            left + Inches(0.1), Inches(1.6),
            Inches(2.7), Inches(0.5)
        )
        text_box.text_frame.text = comp_title
        text_box.text_frame.paragraphs[0].font.size = Pt(16)
        text_box.text_frame.paragraphs[0].font.bold = True
        text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add description
        desc_box = slide.shapes.add_textbox(
            left + Inches(0.2), Inches(2.3),
            Inches(2.5), Inches(2)
        )
        desc_box.text_frame.text = comp_text
        desc_box.text_frame.paragraphs[0].font.size = Pt(12)
        desc_box.text_frame.paragraphs[0].font.color.rgb = text_dark
    
    # Slide 4: Container Environment - Technical specs
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_content_slide(slide, accent_teal)
    
    title = slide.shapes.title
    title.text = "🐳 Sandboxed Container Environment"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Technical Specifications section
    p = tf.add_paragraph()
    p.text = "Technical Specifications"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = primary_blue
    p.space_after = Pt(8)
    
    specs = [
        "Ubuntu 24.04 LTS base image",
        "ZSH shell for command execution",
        "Python 3.12 with data science libraries",
        "Node.js 18.x for JavaScript operations",
        "Persistent /mnt directory across session"
    ]
    
    for spec in specs:
        p = tf.add_paragraph()
        p.text = f"• {spec}"
        p.font.size = Pt(14)
        p.font.color.rgb = text_dark
        p.level = 1
        p.space_after = Pt(6)
    
    # Pre-installed Tools section
    p = tf.add_paragraph()
    p.text = "\nPre-installed Tools"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = accent_teal
    p.space_after = Pt(8)
    
    tools = [
        "Document processing: pandoc, LibreOffice",
        "PowerPoint creation: python-pptx, pptxgenjs",
        "Image manipulation: ImageMagick",
        "Development tools: git, compilers"
    ]
    
    for tool in tools:
        p = tf.add_paragraph()
        p.text = f"• {tool}"
        p.font.size = Pt(14)
        p.font.color.rgb = text_dark
        p.level = 1
        p.space_after = Pt(6)
    
    # Slide 5: Tool Capabilities - Grid layout
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    style_blank_slide(slide, primary_blue)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🔧 Integrated Tool Capabilities"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = dark_blue
    
    # Create capability cards
    capabilities = [
        ("📁 File Management", ["View, create, edit text files", "Navigate directories", "Image processing"], Inches(0.5), Inches(1.3)),
        ("⚡ System Operations", ["Execute shell commands", "Container monitoring", "Environment reset"], Inches(3.5), Inches(1.3)),
        ("🌐 External Integration", ["Web search & fetching", "Email sending", "Calendar access"], Inches(6.5), Inches(1.3)),
        ("📊 Analysis Tools", ["Codebase analysis", "Data processing", "Document conversion"], Inches(2), Inches(3.3))
    ]
    
    for title_text, items, left, top in capabilities:
        # Card background
        card = add_rounded_rectangle(
            slide, left, top,
            Inches(2.8), Inches(1.8),
            RGBColor(255, 255, 255)
        )
        card.shadow.visible = True
        card.shadow.blur_radius = Pt(8)
        
        # Card title
        title_box = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.1),
            Inches(2.6), Inches(0.4)
        )
        title_box.text_frame.text = title_text
        title_box.text_frame.paragraphs[0].font.size = Pt(14)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = primary_blue
        
        # Card items
        items_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.5),
            Inches(2.4), Inches(1.2)
        )
        for item in items:
            p = items_box.text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = text_light
    
    # Slide 6: Security Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_content_slide(slide, RGBColor(239, 68, 68))  # Red accent for security
    
    title = slide.shapes.title
    title.text = "🔐 Security & Isolation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("🛡️ Container Isolation", [
            "Sandboxed Docker environment",
            "No direct host filesystem access",
            "Resource limits and constraints",
            "Separate workspace per conversation"
        ]),
        ("🎯 Controlled Host Access", [
            "Specific tools for host operations",
            "Explicit permission boundaries",
            "Read-only codebase analysis"
        ]),
        ("🔒 Data Protection", [
            "Session-based isolation",
            "No cross-conversation data access"
        ])
    ]
    
    for section_title, items in sections:
        p = tf.add_paragraph()
        p.text = section_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = dark_blue
        p.space_after = Pt(6)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = text_dark
            p.level = 1
            p.space_after = Pt(4)
        
        p.space_after = Pt(10)
    
    # Slide 7: Processing Pipeline - Flow diagram style
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    style_blank_slide(slide, accent_purple)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🔄 Request Processing Pipeline"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = dark_blue
    
    # Create pipeline stages
    stages = [
        ("1️⃣ Input Processing", ["Natural language parsing", "Intent recognition", "Context integration"]),
        ("2️⃣ Planning & Reasoning", ["Task decomposition", "Tool selection logic", "Parameter extraction"]),
        ("3️⃣ Execution", ["Function invocation", "Result processing", "Error handling"]),
        ("4️⃣ Response Generation", ["Result synthesis", "Natural language formatting", "Output delivery"])
    ]
    
    for i, (stage_title, items) in enumerate(stages):
        left = Inches(0.5 + i * 2.3)
        top = Inches(1.5)
        
        # Stage box
        box = add_rounded_rectangle(
            slide, left, top,
            Inches(2.1), Inches(3),
            RGBColor(255, 255, 255)
        )
        box.shadow.visible = True
        
        # Stage header with gradient effect
        header = add_rounded_rectangle(
            slide, left, top,
            Inches(2.1), Inches(0.5),
            [primary_blue, accent_purple][i % 2]
        )
        
        # Stage title
        title_box = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.05),
            Inches(1.9), Inches(0.4)
        )
        title_box.text_frame.text = stage_title
        title_box.text_frame.paragraphs[0].font.size = Pt(13)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Stage items
        items_box = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.7),
            Inches(1.8), Inches(2)
        )
        for item in items:
            p = items_box.text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = text_dark
            p.space_after = Pt(8)
        
        # Add arrow between stages (except for last)
        if i < len(stages) - 1:
            arrow_box = slide.shapes.add_textbox(
                left + Inches(2.15), top + Inches(1.3),
                Inches(0.1), Inches(0.4)
            )
            arrow_box.text_frame.text = "→"
            arrow_box.text_frame.paragraphs[0].font.size = Pt(24)
            arrow_box.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    # Slide 8: Memory & Context
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_content_slide(slide, accent_teal)
    
    title = slide.shapes.title
    title.text = "🧠 Memory & Context Management"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    memory_sections = [
        ("💭 Conversation Context", [
            "Full conversation history retention",
            "Tool execution results tracking",
            "Error and retry management"
        ]),
        ("👤 User Memory System", [
            "Persistent fact storage",
            "Cross-conversation memory",
            "User preference tracking"
        ]),
        ("📚 Knowledge Base", [
            "Training data up to January 2025",
            "Real-time web search integration",
            "Document analysis capabilities"
        ])
    ]
    
    for section_title, items in memory_sections:
        p = tf.add_paragraph()
        p.text = section_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = accent_teal
        p.space_after = Pt(6)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = text_dark
            p.level = 1
            p.space_after = Pt(4)
        
        p.space_after = Pt(10)
    
    # Slide 9: Technical Stack
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_content_slide(slide, primary_blue)
    
    title = slide.shapes.title
    title.text = "💻 Technical Stack Details"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    stack_sections = [
        ("🔤 Programming Languages", [
            "Python 3.12 - Primary scripting",
            "JavaScript/Node.js - Web tools",
            "Shell scripting - System operations"
        ]),
        ("📦 Key Libraries", [
            "python-pptx - Presentation creation",
            "pandas - Data manipulation",
            "Pillow - Image processing",
            "markitdown - Document conversion"
        ]),
        ("🏗️ Infrastructure", [
            "Docker containerization",
            "Linux-based execution environment",
            "RESTful API interfaces"
        ])
    ]
    
    for section_title, items in stack_sections:
        p = tf.add_paragraph()
        p.text = section_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = primary_blue
        p.space_after = Pt(6)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = text_dark
            p.level = 1
            p.space_after = Pt(4)
        
        p.space_after = Pt(10)
    
    # Slide 10: Performance
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_content_slide(slide, accent_purple)
    
    title = slide.shapes.title
    title.text = "⚡ Performance Characteristics"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    perf_sections = [
        ("🚀 Response Generation", [
            "Context-aware processing",
            "Streaming response capability",
            "Multi-turn conversation support"
        ]),
        ("⏱️ Execution Limits", [
            "300-second default command timeout",
            "10-minute codebase analysis timeout",
            "File size handling up to 10MB"
        ]),
        ("⚙️ Optimization Features", [
            "Parallel tool execution",
            "Result caching within session",
            "Efficient memory management"
        ])
    ]
    
    for section_title, items in perf_sections:
        p = tf.add_paragraph()
        p.text = section_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = accent_purple
        p.space_after = Pt(6)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = text_dark
            p.level = 1
            p.space_after = Pt(4)
        
        p.space_after = Pt(10)
    
    # Slide 11: Use Cases with icons
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    style_blank_slide(slide, accent_teal)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🎯 Primary Use Cases"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = dark_blue
    
    # Create use case cards in 2x2 grid
    use_cases = [
        ("📝 Document Creation", ["PowerPoint presentations", "Technical reports", "Data analysis docs"], Inches(0.8), Inches(1.4)),
        ("💻 Development Support", ["Code writing & debugging", "Codebase analysis", "Script automation"], Inches(5.2), Inches(1.4)),
        ("📧 Productivity Tasks", ["Email composition", "Calendar management", "File organization"], Inches(0.8), Inches(3.4)),
        ("🔍 Information Processing", ["Web research", "Data extraction", "Content analysis"], Inches(5.2), Inches(3.4))
    ]
    
    for title_text, items, left, top in use_cases:
        # Card
        card = add_rounded_rectangle(
            slide, left, top,
            Inches(3.8), Inches(1.6),
            RGBColor(255, 255, 255)
        )
        card.shadow.visible = True
        
        # Accent strip
        accent = add_rounded_rectangle(
            slide, left, top,
            Inches(0.08), Inches(1.6),
            accent_teal
        )
        
        # Title
        title_box = slide.shapes.add_textbox(
            left + Inches(0.3), top + Inches(0.15),
            Inches(3.3), Inches(0.4)
        )
        title_box.text_frame.text = title_text
        title_box.text_frame.paragraphs[0].font.size = Pt(16)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = dark_blue
        
        # Items
        items_box = slide.shapes.add_textbox(
            left + Inches(0.4), top + Inches(0.55),
            Inches(3.2), Inches(0.9)
        )
        for item in items:
            p = items_box.text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = text_light
    
    # Slide 12: Architecture Advantages
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_content_slide(slide, primary_blue)
    
    title = slide.shapes.title
    title.text = "✨ Architecture Advantages"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    advantages = [
        ("🔄 Flexibility", [
            "Extensible tool framework",
            "Multiple language support",
            "Adaptable to various workflows"
        ]),
        ("🛡️ Reliability", [
            "Isolated execution environment",
            "Error recovery mechanisms",
            "Consistent performance"
        ]),
        ("📈 Scalability", [
            "Container-based architecture",
            "Modular component design",
            "Efficient resource utilization"
        ]),
        ("👥 User Experience", [
            "Natural language interface",
            "Context-aware responses",
            "Seamless tool integration"
        ])
    ]
    
    for adv_title, items in advantages:
        p = tf.add_paragraph()
        p.text = adv_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = primary_blue
        p.space_after = Pt(6)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = text_dark
            p.level = 1
            p.space_after = Pt(4)
        
        p.space_after = Pt(10)
    
    # Slide 13: Summary
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Special styling for summary
    bg_shape = add_rounded_rectangle(
        slide,
        Inches(0), Inches(0),
        Inches(10), Inches(5.625),
        dark_blue
    )
    slide.shapes._spTree.remove(bg_shape._element)
    slide.shapes._spTree.insert(2, bg_shape._element)
    
    title = slide.shapes.title
    title.text = "🎯 Summary"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    summary_text = [
        "Stella represents a sophisticated AI assistant architecture that:",
        "",
        "✅ Combines advanced language understanding with practical tool execution",
        "",
        "✅ Maintains security through containerized isolation",
        "",
        "✅ Provides comprehensive productivity capabilities",
        "",
        "✅ Integrates seamlessly with user workflows",
        "",
        "✅ Delivers consistent, reliable performance",
        "",
        "Built on Claude's foundation with custom enhancements by Jess Solutions"
    ]
    
    for line in summary_text:
        p = tf.add_paragraph()
        p.text = line
        if line.startswith("✅"):
            p.font.size = Pt(18)
            p.font.color.rgb = secondary_blue
            p.font.bold = True
        elif line.startswith("Stella"):
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
        elif line.startswith("Built on"):
            p.font.size = Pt(14)
            p.font.color.rgb = secondary_blue
            p.font.italic = True
        else:
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(8)
    
    # Save presentation
    prs.save('conversation_data/stella_architecture_aesthetic.pptx')
    print("Aesthetic presentation created successfully!")

# Create the presentation
create_aesthetic_stella_presentation()
