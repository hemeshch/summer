from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "The Fascinating World of Fish"
subtitle.text = "A Deep Dive into Aquatic Life"

# Apply ocean theme colors to title slide
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 48, 135)
title.text_frame.paragraphs[0].font.size = Pt(44)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 119, 200)

# Slide 2: Introduction
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
content2 = slide2.placeholders[1]
title2.text = "What Makes Fish Special?"
content2.text = "• Over 34,000 species of fish exist worldwide\n\n• Fish were the first vertebrates with backbones\n\n• They've existed for more than 500 million years\n\n• Fish can be found in nearly every aquatic environment\n\n• From tiny gobies (8mm) to massive whale sharks (12m)"

# Slide 3: Amazing Adaptations
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
content3 = slide3.placeholders[1]
title3.text = "Incredible Fish Adaptations"
content3.text = "🐟 Lateral Line System\n• Detects vibrations and water movement\n\n🫧 Swim Bladders\n• Control buoyancy without constant swimming\n\n👁️ 360° Vision\n• Most fish can see in all directions simultaneously\n\n🌡️ Cold-Blooded Masters\n• Body temperature matches their environment"

# Slide 4: Bioluminescence
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
content4 = slide4.placeholders[1]
title4.text = "Deep Sea Light Shows"
content4.text = "The Anglerfish\n• Uses a glowing lure to attract prey\n• Lives in complete darkness up to 1,500m deep\n\nLanternfish\n• Most abundant vertebrates on Earth\n• Create their own light for communication\n\nFlashlight Fish\n• Have light organs under their eyes\n• Can 'turn off' their lights by covering them with skin"

# Slide 5: Speed Demons
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
content5 = slide5.placeholders[1]
title5.text = "The Ocean's Speed Champions"
content5.text = "🥇 Black Marlin\n• Top speed: 82 mph (132 km/h)\n\n🥈 Sailfish\n• Top speed: 68 mph (110 km/h)\n• Distinctive sail-like dorsal fin\n\n🥉 Swordfish\n• Top speed: 60 mph (97 km/h)\n• Uses its 'sword' to slash prey\n\n⚡ Mako Shark\n• Top speed: 46 mph (74 km/h)\n• Can leap 20 feet out of water"

# Slide 6: Weird and Wonderful
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title6 = slide6.shapes.title
content6 = slide6.placeholders[1]
title6.text = "Nature's Most Unusual Fish"
content6.text = "Blobfish\n• Gelatinous body with no muscles\n• Lives at depths where pressure is 120x surface level\n\nLeafy Sea Dragon\n• Master of camouflage\n• Looks exactly like floating seaweed\n\nMandarin Fish\n• No scales, covered in toxic mucus\n• One of only two vertebrates with blue pigmentation\n\nBoxfish\n• Cube-shaped body armor\n• Releases deadly toxins when threatened"

# Slide 7: Fish Intelligence
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
title7 = slide7.shapes.title
content7 = slide7.placeholders[1]
title7.text = "Smarter Than You Think!"
content7.text = "🧠 Problem Solving\n• Archerfish calculate physics to shoot water at insects\n• Tuskfish use rocks as tools to crack shells\n\n🎯 Memory\n• Goldfish can remember things for at least 3 months\n• Some fish recognize human faces\n\n👥 Social Intelligence\n• Cleaner fish run 'cleaning stations' with regular clients\n• Groupers team up with eels for cooperative hunting\n\n🗺️ Navigation\n• Salmon return to exact birthplace after years at sea"

# Slide 8: Conservation
slide8 = prs.slides.add_slide(prs.slide_layouts[1])
title8 = slide8.shapes.title
content8 = slide8.placeholders[1]
title8.text = "Protecting Our Finned Friends"
content8.text = "Current Threats:\n• Overfishing affects 90% of fish stocks\n• Ocean acidification from climate change\n• Plastic pollution kills millions annually\n• Habitat destruction in coral reefs and rivers\n\nHow You Can Help:\n• Choose sustainable seafood\n• Reduce plastic use\n• Support marine protected areas\n• Learn and share knowledge about fish"

# Slide 9: Fun Facts
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
title9 = slide9.shapes.title
content9 = slide9.placeholders[1]
title9.text = "Mind-Blowing Fish Facts"
content9.text = "🎵 Fish can be noisy! Drums, croakers, and grunts are named after their sounds\n\n😴 Parrotfish sleep in mucus bubbles for protection\n\n🔄 Some fish can change gender during their lifetime\n\n❄️ Antarctic icefish have antifreeze proteins in their blood\n\n👑 A group of goldfish is called a 'troubling'\n\n💪 Electric eels can generate 600 volts of electricity\n\n🎨 Cuttlefish can create moving patterns on their skin like LED displays"

# Slide 10: Conclusion
slide10 = prs.slides.add_slide(prs.slide_layouts[1])
title10 = slide10.shapes.title
content10 = slide10.placeholders[1]
title10.text = "The Ocean's Endless Wonders"
content10.text = "Fish have survived mass extinctions, evolved incredible abilities, and continue to surprise scientists with new discoveries.\n\nFrom the smallest reef dwellers to the giants of the deep, fish showcase nature's creativity and resilience.\n\nNext time you see a fish, remember:\nYou're looking at a member of Earth's most successful and diverse group of vertebrates!\n\n🐠 Thank you for diving deep with us! 🐠"

# Save presentation
prs.save('/mnt/conversation_data/cool_fish_presentation.pptx')
print("Fish presentation created successfully!")
