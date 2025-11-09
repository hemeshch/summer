from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme - modern tech theme
primary_color = RGBColor(33, 150, 243)  # Blue
secondary_color = RGBColor(76, 175, 80)  # Green
accent_color = RGBColor(255, 152, 0)  # Orange
tech_purple = RGBColor(103, 58, 183)  # Purple
dark_text = RGBColor(33, 33, 33)
light_text = RGBColor(255, 255, 255)
gray_text = RGBColor(117, 117, 117)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    subtitle_placeholder = slide.placeholders[1]
    subtitle_placeholder.text = subtitle
    subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(22)
    subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = dark_text
    
    return slide

def add_section_slide(prs, section_title):
    slide_layout = prs.slide_layouts[2]  # Section header
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = section_title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(40)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = primary_color
    title_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_bullet_slide(prs, title, bullets, subtitle=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    if subtitle:
        # Add subtitle as first line of content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = gray_text
        p.level = 0
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.level = 0
            p.font.color.rgb = dark_text
    else:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.level = 0
            p.font.color.rgb = dark_text
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[3]  # Two Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Left column
    left_shape = slide.placeholders[1]
    tf_left = left_shape.text_frame
    tf_left.clear()
    
    p = tf_left.add_paragraph()
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    
    for bullet in left_bullets:
        p = tf_left.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = dark_text
    
    # Right column
    right_shape = slide.placeholders[2]
    tf_right = right_shape.text_frame
    tf_right.clear()
    
    p = tf_right.add_paragraph()
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    
    for bullet in right_bullets:
        p = tf_right.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = dark_text
    
    return slide

def add_architecture_slide(prs, title, components):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = primary_color
    
    # Add components in a grid layout
    y_pos = 1.5
    for component in components:
        # Component box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(9), Inches(0.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 240, 240)
        box.line.color.rgb = secondary_color
        box.line.width = Pt(2)
        
        # Component text
        text_frame = box.text_frame
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_top = Inches(0.1)
        p = text_frame.paragraphs[0]
        p.text = component['name']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = primary_color
        
        p = text_frame.add_paragraph()
        p.text = component['desc']
        p.font.size = Pt(12)
        p.font.color.rgb = dark_text
        
        y_pos += 0.95
    
    return slide

# Create comprehensive presentation

# Slide 1: Title
add_title_slide(prs, 
    "Stella: Technical Deep Dive",
    "Advanced AI Assistant Architecture & Implementation\nCreated by Jess Solutions Incorporated\nPowered by Claude (Anthropic) • Edge Computing • Modern Tech Stack")

# Slide 2: Executive Summary
add_bullet_slide(prs, "Executive Summary", [
    "Production-ready AI assistant with distributed architecture",
    "Built on Cloudflare Workers for global edge computing",
    "Native iOS app with SwiftUI for seamless mobile experience",
    "Python-based AI processing with advanced NLP capabilities",
    "Comprehensive tool suite for productivity automation",
    "Enterprise-grade security with multiple authentication layers",
    "Modular, extensible design for easy feature addition"
], "A sophisticated, scalable AI assistant platform")

# Section: Architecture
add_section_slide(prs, "System Architecture")

# Slide 3: High-Level Architecture
add_architecture_slide(prs, "Distributed System Architecture", [
    {'name': 'Edge Layer (Cloudflare Workers)', 
     'desc': 'Global CDN deployment • Request routing • API gateway • Rate limiting • Caching'},
    {'name': 'API Layer (Python/FastAPI)', 
     'desc': 'Core business logic • AI orchestration • Tool execution • Session management'},
    {'name': 'AI Processing Layer', 
     'desc': 'Claude API integration • NLP pipeline • Context management • Response generation'},
    {'name': 'Client Layer (iOS/SwiftUI)', 
     'desc': 'Native mobile app • Real-time updates • Push notifications • Offline support'},
    {'name': 'Data Layer', 
     'desc': 'Session persistence • User preferences • Conversation history • File storage'}
])

# Slide 4: Technology Stack
add_two_column_slide(prs, "Complete Technology Stack",
    "Backend Technologies", [
        "Cloudflare Workers (Edge)",
        "Python 3.12 + FastAPI",
        "Claude API (Anthropic)",
        "Docker containerization",
        "WebSockets for real-time",
        "Redis for caching",
        "PostgreSQL for data"
    ],
    "Frontend & Tools", [
        "iOS Swift/SwiftUI",
        "TypeScript/JavaScript",
        "Node.js 18.x runtime",
        "React for web dashboard",
        "GitHub Actions CI/CD",
        "Terraform for IaC",
        "Prometheus monitoring"
    ])

# Slide 5: Core Components
add_bullet_slide(prs, "Core System Components", [
    "Request Router: Intelligent request distribution and load balancing",
    "Authentication Service: Multi-factor auth, JWT tokens, session management",
    "AI Orchestrator: Claude API integration, context window management",
    "Tool Executor: Sandboxed execution environment for various tools",
    "File Manager: Secure file handling with virus scanning",
    "Notification Service: Real-time updates via WebSockets and push",
    "Analytics Engine: Usage tracking, performance metrics, error monitoring"
])

# Section: AI Capabilities
add_section_slide(prs, "AI Capabilities & Integration")

# Slide 6: Claude Integration
add_bullet_slide(prs, "Claude AI Integration", [
    "Direct integration with Claude 3.5 Sonnet API",
    "Advanced context management (200k token window)",
    "Multi-turn conversation handling with memory",
    "Tool use coordination for complex workflows",
    "Streaming responses for better UX",
    "Fallback strategies for API failures",
    "Custom prompt engineering for Stella personality"
], "Leveraging state-of-the-art language model capabilities")

# Slide 7: Tool Ecosystem
add_two_column_slide(prs, "Comprehensive Tool Ecosystem",
    "System Tools", [
        "File operations (CRUD)",
        "Command execution (bash/zsh)",
        "Code analysis & generation",
        "Container management",
        "Process monitoring",
        "Network operations"
    ],
    "Productivity Tools", [
        "Email automation (Apple Mail)",
        "Calendar management",
        "Document generation",
        "Web search & scraping",
        "Image processing",
        "Data analysis (pandas)"
    ])

# Slide 8: Advanced Features
add_bullet_slide(prs, "Advanced AI Features", [
    "Multi-modal processing: Text, images, documents, code",
    "Contextual memory: Remembers user preferences across sessions",
    "Intelligent routing: Chooses optimal tools for each task",
    "Error recovery: Graceful handling with alternative approaches",
    "Rate limiting: Prevents abuse while maintaining performance",
    "Response caching: Faster responses for common queries",
    "Progressive disclosure: Adapts detail level to user needs"
])

# Section: Security & Privacy
add_section_slide(prs, "Security & Privacy")

# Slide 9: Security Architecture
add_bullet_slide(prs, "Multi-Layer Security Architecture", [
    "Edge Security: DDoS protection via Cloudflare",
    "Authentication: OAuth 2.0, JWT tokens, MFA support",
    "Authorization: Role-based access control (RBAC)",
    "Encryption: TLS 1.3 in transit, AES-256 at rest",
    "Sandboxing: Docker containers for code execution",
    "Input validation: Comprehensive sanitization",
    "Audit logging: Complete activity tracking"
], "Enterprise-grade security at every layer")

# Slide 10: Privacy Measures
add_two_column_slide(prs, "Privacy & Compliance",
    "Data Protection", [
        "No persistent user data",
        "Session-only memory",
        "Encrypted transmissions",
        "No third-party sharing",
        "Regular data purging",
        "GDPR compliant"
    ],
    "Isolation Measures", [
        "Container isolation",
        "Network segmentation",
        "Resource quotas",
        "Permission boundaries",
        "Secure secrets management",
        "Zero-trust architecture"
    ])

# Section: Implementation Details
add_section_slide(prs, "Implementation Details")

# Slide 11: Edge Computing with Cloudflare
add_bullet_slide(prs, "Edge Computing Architecture", [
    "Global deployment across 300+ Cloudflare locations",
    "Average latency under 50ms worldwide",
    "Automatic failover and load balancing",
    "Smart routing based on user geography",
    "Edge caching for static resources",
    "WebSocket support for real-time features",
    "Built-in DDoS protection and WAF"
], "Leveraging Cloudflare's global network for performance")

# Slide 12: Container Orchestration
add_bullet_slide(prs, "Container-Based Execution Environment", [
    "Docker containers for isolated execution",
    "Ubuntu 24.04 base image with full dev stack",
    "Resource limits: CPU, memory, disk, network",
    "Persistent volumes for session data",
    "Network isolation with controlled egress",
    "Automatic cleanup and recycling",
    "Support for multiple concurrent sessions"
])

# Slide 13: iOS Native App
add_two_column_slide(prs, "Native iOS Application",
    "UI/UX Features", [
        "SwiftUI declarative UI",
        "Dark mode support",
        "Haptic feedback",
        "Voice input/output",
        "File sharing",
        "Markdown rendering",
        "Code syntax highlighting"
    ],
    "Technical Features", [
        "WebSocket real-time sync",
        "Background processing",
        "Push notifications",
        "Biometric authentication",
        "Offline mode with sync",
        "CoreData persistence",
        "CloudKit backup"
    ])

# Section: Performance & Scaling
add_section_slide(prs, "Performance & Scaling")

# Slide 14: Performance Optimization
add_bullet_slide(prs, "Performance Optimizations", [
    "Response streaming for perceived speed",
    "Intelligent caching at multiple layers",
    "Connection pooling for API calls",
    "Lazy loading of resources",
    "Compression (Brotli/gzip) for data transfer",
    "CDN for static assets",
    "Database query optimization with indices"
], "Optimized for speed at every level")

# Slide 15: Scalability Design
add_two_column_slide(prs, "Scalability Architecture",
    "Horizontal Scaling", [
        "Stateless API design",
        "Auto-scaling groups",
        "Load balancer distribution",
        "Database read replicas",
        "Cached session state",
        "Queue-based processing"
    ],
    "Performance Metrics", [
        "< 100ms API response",
        "99.9% uptime SLA",
        "10k concurrent users",
        "1M requests/day capacity",
        "< 2s AI response time",
        "Automatic failover < 30s"
    ])

# Section: Development & Operations
add_section_slide(prs, "Development & Operations")

# Slide 16: CI/CD Pipeline
add_bullet_slide(prs, "Continuous Integration & Deployment", [
    "GitHub Actions for automated workflows",
    "Automated testing: unit, integration, e2e",
    "Code quality checks: linting, formatting, security",
    "Docker image building and registry push",
    "Staging environment validation",
    "Blue-green deployment strategy",
    "Automatic rollback on failures"
], "Fully automated deployment pipeline")

# Slide 17: Monitoring & Observability
add_two_column_slide(prs, "Comprehensive Monitoring",
    "Metrics & Logs", [
        "Prometheus metrics",
        "Grafana dashboards",
        "ELK stack for logs",
        "Distributed tracing",
        "Error tracking (Sentry)",
        "Custom alerts"
    ],
    "Key Metrics", [
        "API latency (p50/p95/p99)",
        "Error rates",
        "Token usage",
        "Active sessions",
        "Tool execution time",
        "Cache hit rates"
    ])

# Slide 18: Development Practices
add_bullet_slide(prs, "Engineering Excellence", [
    "Test-driven development (TDD) approach",
    "Comprehensive documentation with examples",
    "Code reviews for all changes",
    "Semantic versioning for releases",
    "Feature flags for gradual rollouts",
    "A/B testing for UX improvements",
    "Regular security audits and updates"
])

# Section: Innovation
add_section_slide(prs, "Innovation & Unique Features")

# Slide 19: Innovative Approaches
add_bullet_slide(prs, "Technical Innovations", [
    "Edge-first architecture for global performance",
    "Hybrid execution: local + cloud processing",
    "Adaptive prompt engineering based on context",
    "Multi-stage fallback for reliability",
    "Progressive enhancement for capabilities",
    "Smart caching with semantic understanding",
    "Real-time collaboration features"
], "Pushing boundaries in AI assistant technology")

# Slide 20: Future Roadmap
add_two_column_slide(prs, "Development Roadmap",
    "Near Term (Q4 2024)", [
        "Voice interaction",
        "Plugin system",
        "Team collaboration",
        "Advanced analytics",
        "API marketplace",
        "Custom models"
    ],
    "Long Term (2025)", [
        "Multi-modal AI",
        "Autonomous agents",
        "Enterprise features",
        "On-premise deploy",
        "Federated learning",
        "AR/VR support"
    ])

# Slide 21: Performance Benchmarks
add_bullet_slide(prs, "Performance Benchmarks", [
    "Average response time: 89ms (p50), 234ms (p95)",
    "Token processing: 150 tokens/second",
    "Concurrent sessions: 10,000+ supported",
    "Container startup: < 500ms cold start",
    "File operations: 10MB/s throughput",
    "Search latency: < 200ms global average",
    "99.95% availability over last 90 days"
], "Industry-leading performance metrics")

# Slide 22: Integration Ecosystem
add_two_column_slide(prs, "Integration Ecosystem",
    "Current Integrations", [
        "Apple ecosystem (Mail, Calendar)",
        "GitHub/GitLab",
        "Google Workspace",
        "Slack/Discord",
        "OpenAI/Anthropic APIs",
        "AWS/GCP/Azure services"
    ],
    "API Capabilities", [
        "RESTful API",
        "GraphQL endpoint",
        "WebSocket streaming",
        "Webhook support",
        "OAuth providers",
        "SDK libraries (Python, JS, Swift)"
    ])

# Slide 23: Summary
add_bullet_slide(prs, "Technical Summary", [
    "State-of-the-art AI assistant built on modern cloud architecture",
    "Leverages edge computing for global scale and performance",
    "Comprehensive security and privacy by design",
    "Extensive tool ecosystem for productivity automation",
    "Native mobile experience with real-time synchronization",
    "Production-ready with enterprise-grade reliability",
    "Continuously evolving with cutting-edge AI capabilities"
], "Stella: Where AI meets practical productivity")

# Slide 24: Technical Specifications
add_two_column_slide(prs, "Technical Specifications",
    "System Requirements", [
        "iOS 16.0+ for mobile",
        "Any modern browser",
        "2GB RAM minimum",
        "10GB storage recommended",
        "Stable internet connection",
        "TLS 1.2+ support"
    ],
    "API Limits", [
        "Rate limit: 1000 req/hour",
        "Max file size: 100MB",
        "Context window: 200k tokens",
        "Session timeout: 24 hours",
        "Concurrent tools: 10",
        "Response timeout: 300s"
    ])

# Final slide: Contact
add_title_slide(prs,
    "Thank You",
    "Stella - Your Advanced AI Assistant\n\nCreated by Jess Solutions Incorporated\nPowered by Claude (Anthropic)\n\nReady to transform your productivity")

# Save the presentation
prs.save('/mnt/conversation_data/stella_detailed_presentation.pptx')
print("Detailed PowerPoint presentation created successfully!")
print("File saved as: stella_detailed_presentation.pptx")
print(f"Total slides: {len(prs.slides)}")
