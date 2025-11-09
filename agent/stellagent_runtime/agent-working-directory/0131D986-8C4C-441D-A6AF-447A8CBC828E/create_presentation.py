from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme - modern tech theme
primary_color = RGBColor(33, 150, 243)  # Blue
secondary_color = RGBColor(76, 175, 80)  # Green
accent_color = RGBColor(255, 152, 0)  # Orange
dark_text = RGBColor(33, 33, 33)
light_text = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Subtitle
    subtitle_placeholder = slide.placeholders[1]
    subtitle_placeholder.text = subtitle
    subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = dark_text
    
    return slide

def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.level = 0
        p.font.color.rgb = dark_text
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[3]  # Two Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Left column
    left_shape = slide.placeholders[1]
    tf_left = left_shape.text_frame
    tf_left.clear()
    
    # Add left title
    p = tf_left.add_paragraph()
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    
    for bullet in left_bullets:
        p = tf_left.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = dark_text
    
    # Right column
    right_shape = slide.placeholders[2]
    tf_right = right_shape.text_frame
    tf_right.clear()
    
    # Add right title
    p = tf_right.add_paragraph()
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    
    for bullet in right_bullets:
        p = tf_right.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = dark_text
    
    return slide

# Create slides

# Slide 1: Title
add_title_slide(prs, 
    "Stella",
    "Your Productivity-Focused AI Assistant\nCreated by Jess Solutions Incorporated\nPowered by Claude Architecture (Anthropic)")

# Slide 2: Core Identity
add_bullet_slide(prs, "Core Identity", [
    "Professional AI assistant specialized in productivity and digital task management",
    "Maintains helpful, professional, and friendly communication style",
    "Transparent about AI nature while focusing on practical capabilities",
    "Knowledge base current through January 2025",
    "Tailored responses based on task complexity and user needs"
])

# Slide 3: Technical Foundation
add_bullet_slide(prs, "Technical Foundation", [
    "Built on Claude architecture developed by Anthropic",
    "Advanced natural language processing and understanding",
    "Multi-modal capabilities for text and image analysis",
    "Operates within sandboxed Linux container environment",
    "Ubuntu 24.04 base with comprehensive development tools"
])

# Slide 4: Primary Capabilities
add_bullet_slide(prs, "Primary Capabilities", [
    "File Management & Organization",
    "Code Execution & Development Workflows",
    "Email & Calendar Management (Apple ecosystem)",
    "Document Creation & Editing",
    "System Automation & Task Execution",
    "Web Search & Information Retrieval",
    "Codebase Analysis & Documentation"
])

# Slide 5: Technical Environment
add_two_column_slide(prs, "Technical Environment",
    "Core System", [
        "Docker Container (Isolated)",
        "Ubuntu 24.04 LTS",
        "ZSH Shell",
        "Persistent workspace at /mnt"
    ],
    "Development Stack", [
        "Python 3.12 + Data Science libs",
        "Node.js 18.x",
        "PowerPoint tools suite",
        "Document processing tools"
    ])

# Slide 6: Tool Suite Overview
add_two_column_slide(prs, "Comprehensive Tool Suite",
    "Container Operations", [
        "Command execution",
        "File management",
        "Code development",
        "Data processing",
        "Document generation"
    ],
    "Host Machine Access", [
        "Calendar integration",
        "Email automation",
        "Codebase analysis",
        "Web content fetching",
        "User preference storage"
    ])

# Slide 7: Communication Approach
add_bullet_slide(prs, "Communication Philosophy", [
    "Concise responses for simple queries",
    "Thorough explanations for complex challenges",
    "Clear, structured formatting for technical instructions",
    "Natural, conversational tone for discussions",
    "Objective and honest feedback, even when difficult",
    "Practical examples to illustrate concepts"
])

# Slide 8: Security & Privacy
add_bullet_slide(prs, "Security & Privacy", [
    "Sandboxed execution environment prevents system damage",
    "Isolated from host system (except for specific authorized tools)",
    "Transparent about capabilities and limitations",
    "Respects intellectual property and copyright",
    "User data remains within conversation context",
    "No unauthorized access to personal files"
])

# Slide 9: Key Strengths
add_bullet_slide(prs, "Key Strengths", [
    "Multi-tool orchestration for complex workflows",
    "Clear explanation of technical concepts with examples",
    "Professional document and presentation creation",
    "Code understanding, generation, and debugging",
    "Productivity workflow optimization",
    "Real-time information retrieval and processing"
])

# Slide 10: Use Cases
add_two_column_slide(prs, "Ideal Use Cases",
    "Development & Technical", [
        "Code review and debugging",
        "Script automation",
        "Technical documentation",
        "Data analysis",
        "System administration"
    ],
    "Productivity & Business", [
        "Email drafting",
        "Calendar management",
        "Report generation",
        "Presentation creation",
        "File organization"
    ])

# Slide 11: Summary
add_bullet_slide(prs, "Summary", [
    "Created by Jess Solutions Incorporated",
    "Powered by Claude architecture from Anthropic",
    "Designed to enhance user productivity",
    "Comprehensive tool suite for digital task management",
    "Professional, helpful, and transparent AI assistant",
    "Ready to tackle your productivity challenges"
])

# Save the presentation
prs.save('/mnt/conversation_data/stella_presentation.pptx')
print("PowerPoint presentation created successfully!")
print("File saved as: stella_presentation.pptx")
