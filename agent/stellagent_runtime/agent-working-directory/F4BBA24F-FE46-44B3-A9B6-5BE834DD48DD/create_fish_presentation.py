from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "The World of Fish"
subtitle.text = "An Introduction to Aquatic Life"

# Slide 2: What Are Fish?
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide2.shapes.title
content = slide2.placeholders[1]
title.text = "What Are Fish?"
content.text = (
    "• Cold-blooded vertebrates that live in water\n"
    "• Breathe through gills\n"
    "• Have fins for swimming\n"
    "• Most have scales covering their bodies\n"
    "• Over 34,000 known species worldwide"
)

# Slide 3: Types of Fish
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide3.shapes.title
content = slide3.placeholders[1]
title.text = "Main Types of Fish"
content.text = (
    "Bony Fish (Osteichthyes)\n"
    "• Make up 96% of all fish species\n"
    "• Examples: Salmon, Tuna, Goldfish\n\n"
    "Cartilaginous Fish (Chondrichthyes)\n"
    "• Skeletons made of cartilage\n"
    "• Examples: Sharks, Rays, Skates\n\n"
    "Jawless Fish (Agnatha)\n"
    "• Most primitive type\n"
    "• Examples: Lampreys, Hagfish"
)

# Slide 4: Where Fish Live
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide4.shapes.title
content = slide4.placeholders[1]
title.text = "Fish Habitats"
content.text = (
    "Freshwater Environments\n"
    "• Rivers and streams\n"
    "• Lakes and ponds\n"
    "• About 41% of fish species\n\n"
    "Saltwater Environments\n"
    "• Oceans\n"
    "• Coral reefs\n"
    "• Deep sea trenches\n"
    "• About 58% of fish species\n\n"
    "Brackish Water\n"
    "• Where freshwater meets saltwater\n"
    "• Estuaries and mangroves"
)

# Slide 5: Fish Anatomy
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide5.shapes.title
content = slide5.placeholders[1]
title.text = "Basic Fish Anatomy"
content.text = (
    "External Features:\n"
    "• Fins - for movement and stability\n"
    "• Scales - for protection\n"
    "• Lateral line - detects vibrations\n"
    "• Eyes - vision adapted to water\n\n"
    "Internal Features:\n"
    "• Gills - extract oxygen from water\n"
    "• Swim bladder - controls buoyancy\n"
    "• Heart - two-chambered pump\n"
    "• Simple digestive system"
)

# Slide 6: Interesting Fish Facts
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide6.shapes.title
content = slide6.placeholders[1]
title.text = "Amazing Fish Facts"
content.text = (
    "• The whale shark is the largest fish (up to 40 feet long)\n\n"
    "• The dwarf pygmy goby is the smallest (8mm long)\n\n"
    "• Some fish can change their gender\n\n"
    "• Electric eels can generate 600 volts\n\n"
    "• Fish have been on Earth for over 500 million years\n\n"
    "• Some deep-sea fish produce their own light\n\n"
    "• A school of fish can contain millions of individuals"
)

# Slide 7: Importance of Fish
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide7.shapes.title
content = slide7.placeholders[1]
title.text = "Why Fish Matter"
content.text = (
    "Ecological Importance:\n"
    "• Key part of aquatic food chains\n"
    "• Help maintain ecosystem balance\n"
    "• Nutrient cycling in water bodies\n\n"
    "Economic Importance:\n"
    "• Food source for billions of people\n"
    "• Commercial and recreational fishing\n"
    "• Aquarium trade\n\n"
    "Scientific Importance:\n"
    "• Medical research\n"
    "• Environmental indicators\n"
    "• Evolutionary studies"
)

# Slide 8: Thank You
slide8 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide8.shapes.title
subtitle = slide8.placeholders[1]
title.text = "Thank You!"
subtitle.text = "Questions?"

# Save the presentation
prs.save('conversation_data/fish_presentation.pptx')
print("Presentation created successfully!")
